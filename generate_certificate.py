import qrcode
from pptx import Presentation
from pptx.util import Inches
import os
import sys
from comtypes.client import CreateObject  # For PDF conversion
import time

def generate_certificate(name, certificate_id, issued_date, template_path, output_path):
    # Include .pdf extension in the QR code link
    netlify_link = f"https://codeclashjec.netlify.app/certificates/{certificate_id}.pdf"
    qr_code_path = f"{certificate_id}_qr.png"

    # Generate QR Code
    qr = qrcode.QRCode(version=1, box_size=10, border=5)
    qr.add_data(netlify_link)
    qr.make(fit=True)
    qr_img = qr.make_image(fill="black", back_color="white")
    qr_img.save(qr_code_path)
    print(f"QR Code generated with link: {netlify_link}")  # Debug log

    # Load PowerPoint template
    prs = Presentation(template_path)

    # Replace placeholders
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # Combine all runs into a single string
                    full_text = "".join(run.text for run in paragraph.runs)

                    # Replace placeholders in the combined text
                    full_text = full_text.replace("{{name}}", name)
                    full_text = full_text.replace("{{certificate_id}}", certificate_id)
                    full_text = full_text.replace("{{issued_date}}", issued_date)

                    # Clear all runs and reapply the updated text with formatting
                    for run in paragraph.runs:
                        run.text = ""  # Clear existing text
                    if paragraph.runs:
                        paragraph.runs[0].text = full_text  # Set the updated text in the first run

            # Replace the QR code placeholder with the QR code image
            if shape.has_text_frame and "{{qr_code}}" in shape.text_frame.text:
                # Get the position and size of the textbox
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Delete the placeholder textbox
                sp = shape._element
                sp.getparent().remove(sp)

                # Add the QR code image in the same position and size
                slide.shapes.add_picture(qr_code_path, left, top, width=width, height=height)

    prs.save(output_path)
    print(f"PPTX saved successfully at: {output_path}")  # Debug log

    # Add a short delay to ensure the file is fully written
    time.sleep(1)

    # Check if the file exists before attempting conversion
    if not os.path.exists(output_path):
        raise FileNotFoundError(f"The PPTX file was not created: {output_path}")

    os.remove(qr_code_path)

    # Convert PPTX to PDF
    pdf_output_path = output_path.replace(".pptx", ".pdf")
    print(f"Converting PPTX to PDF: {pdf_output_path}")  # Debug log
    convert_pptx_to_pdf(output_path, pdf_output_path)

    # Delete the PPTX file after successful PDF conversion
    if os.path.exists(output_path):
        os.remove(output_path)
        print(f"PPTX file deleted: {output_path}")  # Debug log

    return pdf_output_path

def convert_pptx_to_pdf(pptx_path, pdf_path):
    try:
        # Convert paths to absolute paths
        pptx_path = os.path.abspath(pptx_path)
        pdf_path = os.path.abspath(pdf_path)

        print(f"Absolute PPTX path: {pptx_path}")  # Debug log
        print(f"Absolute PDF path: {pdf_path}")  # Debug log

        # Initialize PowerPoint application
        powerpoint = CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1

        # Open the presentation
        presentation = powerpoint.Presentations.Open(pptx_path)
        presentation.SaveAs(pdf_path, 32)  # 32 is the format for PDF
        presentation.Close()
        powerpoint.Quit()

        print(f"PDF successfully created at: {pdf_path}")  # Debug log
    except Exception as e:
        print(f"Error during PDF conversion: {e}")  # Debug log
        raise

if __name__ == "__main__":
    name = sys.argv[1]
    certificate_id = sys.argv[2]  # Firebase-generated certificate ID
    issued_date = sys.argv[3]
    template_path = "template.pptx"
    output_path = f"certificates/certificate_{certificate_id}.pptx"

    os.makedirs("certificates", exist_ok=True)
    pdf_path = generate_certificate(name, certificate_id, issued_date, template_path, output_path)
    print(f"Certificate generated: {pdf_path}")